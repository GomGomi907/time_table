from __future__ import annotations

from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
SCREENSHOT_DIR = BASE_DIR / "presentation-screenshots"
OUTPUT_PATH = BASE_DIR / "time_table_mid_progress_report_2026-04-23.pptx"


NAVY = "102542"
NAVY_2 = "17365D"
IVORY = "F7F4EE"
WHITE = "FFFFFF"
PURPLE = "6B4CE6"
PURPLE_SOFT = "EDE6FF"
PURPLE_LINE = "D7CBFF"
MINT = "78C9B5"
MINT_SOFT = "E5F7F2"
AMBER = "F2B566"
AMBER_SOFT = "FFF2DD"
ROSE = "F38A8A"
SLATE = "4C5A6A"
SLATE_SOFT = "6B7887"
LINE = "DDD9EA"
CARD = "FFFDFC"


def color(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def set_slide_size(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def add_bg_fill(slide, fill_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color(fill_hex)


def add_shape(
    slide,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill_hex: str,
    line_hex: str | None = None,
    line_width: float = 1.0,
    transparency: float = 0.0,
) -> None:
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill_hex)
    shape.fill.transparency = transparency
    if line_hex:
        shape.line.color.rgb = color(line_hex)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: int = 16,
    font_name: str = "Malgun Gothic",
    color_hex: str = NAVY,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin_left: float = 0.06,
    margin_right: float = 0.06,
    margin_top: float = 0.03,
    margin_bottom: float = 0.03,
    line_spacing: float = 1.1,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin_left)
    tf.margin_right = Inches(margin_right)
    tf.margin_top = Inches(margin_top)
    tf.margin_bottom = Inches(margin_bottom)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color(color_hex)


def add_chip(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill_hex: str,
    text_hex: str,
    font_size: int = 10,
    line_hex: str | None = None,
    transparency: float = 0.0,
) -> None:
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill_hex=fill_hex,
        line_hex=line_hex,
        transparency=transparency,
    )
    add_text(
        slide,
        text,
        x,
        y + 0.01,
        w,
        h,
        font_size=font_size,
        color_hex=text_hex,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )


def add_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill_hex: str = CARD,
    line_hex: str = LINE,
    transparency: float = 0.0,
) -> None:
    add_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill_hex=fill_hex,
        line_hex=line_hex,
        line_width=1.0,
        transparency=transparency,
    )


def crop_cover_image(path: Path, width_in: float, height_in: float) -> BytesIO:
    target_px_w = max(1200, int(width_in * 220))
    target_px_h = max(700, int(height_in * 220))
    target_ratio = target_px_w / target_px_h

    image = Image.open(path).convert("RGBA")
    src_w, src_h = image.size
    src_ratio = src_w / src_h

    if src_ratio > target_ratio:
        crop_w = int(src_h * target_ratio)
        left = (src_w - crop_w) // 2
        image = image.crop((left, 0, left + crop_w, src_h))
    else:
        crop_h = int(src_w / target_ratio)
        top = (src_h - crop_h) // 2
        image = image.crop((0, top, src_w, top + crop_h))

    image = image.resize((target_px_w, target_px_h), Image.LANCZOS)

    radius = max(24, min(target_px_w, target_px_h) // 16)
    mask = Image.new("L", (target_px_w, target_px_h), 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle((0, 0, target_px_w, target_px_h), radius=radius, fill=255)

    rounded = Image.new("RGBA", (target_px_w, target_px_h), (255, 255, 255, 0))
    rounded.paste(image, (0, 0), mask)

    border = ImageDraw.Draw(rounded)
    border.rounded_rectangle(
        (1, 1, target_px_w - 2, target_px_h - 2),
        radius=radius,
        outline=(215, 203, 255, 255),
        width=max(3, target_px_w // 320),
    )

    stream = BytesIO()
    rounded.save(stream, format="PNG")
    stream.seek(0)
    return stream


def add_picture_card(
    slide,
    image_path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    label: str,
    label_fill: str = PURPLE_SOFT,
    label_text: str = PURPLE,
    card_fill: str = CARD,
) -> None:
    add_card(slide, x, y, w, h, fill_hex=card_fill)
    add_chip(slide, label, x + 0.18, y + 0.16, 1.2, 0.32, fill_hex=label_fill, text_hex=label_text)
    picture_stream = crop_cover_image(image_path, w - 0.28, h - 0.56)
    slide.shapes.add_picture(picture_stream, Inches(x + 0.14), Inches(y + 0.40), width=Inches(w - 0.28), height=Inches(h - 0.48))


def add_floating_picture(
    slide,
    image_path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    label: str | None = None,
) -> None:
    picture_stream = crop_cover_image(image_path, w, h)
    slide.shapes.add_picture(picture_stream, Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if label:
        add_chip(slide, label, x + 0.16, y + 0.12, 1.15, 0.30, fill_hex=WHITE, text_hex=NAVY, font_size=9)


def decorate_dark(slide) -> None:
    add_bg_fill(slide, NAVY)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 8.5, -0.8, 4.0, 4.0, fill_hex=PURPLE, transparency=0.72)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 11.1, 0.6, 2.2, 2.2, fill_hex=MINT, transparency=0.78)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 9.5, 5.0, 4.8, 2.8, fill_hex=PURPLE_SOFT, transparency=0.83)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, -1.2, 5.6, 3.0, 3.0, fill_hex=PURPLE, transparency=0.86)


def decorate_light(slide) -> None:
    add_bg_fill(slide, IVORY)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 10.4, -0.7, 3.2, 3.2, fill_hex=PURPLE_SOFT, transparency=0.35)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, -0.8, 5.4, 2.8, 2.8, fill_hex=MINT_SOFT, transparency=0.45)
    add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, 11.0, 5.5, 2.7, 2.7, fill_hex=AMBER_SOFT, transparency=0.50)


def add_footer(slide, page_text: str, *, dark: bool = False) -> None:
    color_hex = WHITE if dark else SLATE_SOFT
    add_text(
        slide,
        f"Time Table  |  {page_text}",
        0.82,
        7.05,
        3.3,
        0.22,
        font_size=8,
        color_hex=color_hex,
        margin_left=0,
        margin_right=0,
        margin_top=0,
        margin_bottom=0,
    )


def build_cover(slide) -> None:
    decorate_dark(slide)
    add_chip(slide, "MID PROGRESS REPORT", 0.84, 0.56, 1.95, 0.34, fill_hex=PURPLE, text_hex=WHITE, font_size=9)
    add_text(slide, "Time Table", 0.82, 1.05, 4.2, 0.46, font_size=28, color_hex=WHITE, bold=True, margin_left=0)
    add_text(slide, "LLM 기반 일정 자동 조율 플래너", 0.82, 1.48, 4.8, 0.42, font_size=20, color_hex=PURPLE_SOFT, bold=True, margin_left=0)
    add_text(
        slide,
        "주간 시간표·할 일·목표·집중 흐름을 하나의 워크스페이스로 묶고,\nAI가 재배치 제안을 주는 실행형 서비스",
        0.82,
        2.08,
        4.95,
        0.96,
        font_size=15,
        color_hex=WHITE,
        margin_left=0,
        line_spacing=1.2,
    )
    add_text(slide, "중간 진행상황 보고", 0.82, 3.18, 2.4, 0.24, font_size=10, color_hex=MINT, bold=True, margin_left=0)
    add_text(slide, "2026.04.23 기준 저장소 상태를 기준으로 정리", 0.82, 3.41, 3.2, 0.22, font_size=9, color_hex=PURPLE_SOFT, margin_left=0)

    cards = [
        ("목표", "분산된 일정·할 일·목표를 한 화면으로 모아 AI가 주간 실행 흐름을 조율"),
        ("현재", "로그인, 온보딩, 대시보드, 주간 플래너, 집중 모드와 핵심 액션까지 MVP 연결"),
        ("남은 일", "운영형 OAuth, 온보딩 영속화, 설정 UX, sync 실동작, LLM 품질 고도화"),
    ]
    y = 4.10
    for label, body in cards:
        add_card(slide, 0.82, y, 4.85, 0.80, fill_hex=NAVY_2, line_hex="30507B", transparency=0.08)
        add_chip(slide, label, 1.02, y + 0.16, 0.72, 0.28, fill_hex=PURPLE, text_hex=WHITE, font_size=9)
        add_text(slide, body, 1.84, y + 0.10, 3.65, 0.54, font_size=11, color_hex=WHITE, margin_left=0)
        y += 0.92

    add_floating_picture(slide, SCREENSHOT_DIR / "04-schedule.png", 7.05, 0.90, 5.35, 3.10, label="주간 플래너")
    add_floating_picture(slide, SCREENSHOT_DIR / "03-dashboard.png", 6.85, 4.20, 3.25, 2.00, label="대시보드")
    add_floating_picture(slide, SCREENSHOT_DIR / "05-focus.png", 10.35, 3.80, 2.20, 2.60, label="집중 모드")
    add_footer(slide, "01", dark=True)


def build_goal_slide(slide) -> None:
    decorate_light(slide)
    add_chip(slide, "WHAT WE BUILD", 0.82, 0.52, 1.45, 0.34, fill_hex=PURPLE_SOFT, text_hex=PURPLE, font_size=9)
    add_text(
        slide,
        "만들고자 하는 것은 '계획 앱'이 아니라 '실행을 운영하는 워크스페이스'입니다.",
        0.82,
        0.96,
        10.60,
        0.82,
        font_size=21,
        color_hex=NAVY,
        bold=True,
        margin_left=0,
    )
    add_text(
        slide,
        "일정·할 일·목표·집중을 하나의 흐름으로 연결하고, AI는 자동 변경보다 제안과 조정 보조에 집중합니다.",
        0.82,
        1.72,
        9.8,
        0.28,
        font_size=11,
        color_hex=SLATE,
        margin_left=0,
    )

    add_card(slide, 0.82, 1.92, 4.10, 2.36, fill_hex=NAVY, line_hex=NAVY)
    add_chip(slide, "목표", 1.06, 2.16, 0.72, 0.30, fill_hex=PURPLE, text_hex=WHITE, font_size=9)
    add_text(slide, "LLM을 이용한 일정 자동 조율 플래너", 1.04, 2.56, 3.50, 0.84, font_size=19, color_hex=WHITE, bold=True, margin_left=0)
    add_text(
        slide,
        "사용자 캘린더와 생활 패턴을 읽고,\n주간 플랜과 집중 흐름까지 이어지는 실행 경험을 만드는 것이 핵심입니다.",
        1.04,
        3.34,
        3.60,
        0.86,
        font_size=11,
        color_hex=PURPLE_SOFT,
        margin_left=0,
        line_spacing=1.18,
    )

    flow_x = 5.22
    nodes = [
        ("1", "연동 입력", "Google 일정·할 일 연결"),
        ("2", "생활 패턴", "짧은 온보딩으로 루틴 파악"),
        ("3", "주간 플랜", "시간표와 태스크를 한 주 보드로 조합"),
        ("4", "집중 실행", "지금 할 일과 다음 전환 안내"),
        ("5", "AI 재조정", "충돌·지연 시 제안 중심 재배치"),
    ]
    for index, (_, title, desc) in enumerate(nodes):
        x = flow_x + index * 1.48
        add_card(slide, x, 1.98, 1.28, 1.70, fill_hex=WHITE, line_hex=PURPLE_LINE)
        add_chip(slide, str(index + 1), x + 0.40, 2.14, 0.46, 0.28, fill_hex=PURPLE, text_hex=WHITE, font_size=9)
        add_text(slide, title, x + 0.12, 2.56, 1.00, 0.30, font_size=11, color_hex=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + 0.08, 2.92, 1.10, 0.46, font_size=9, color_hex=SLATE, align=PP_ALIGN.CENTER, line_spacing=1.1)
        if index < len(nodes) - 1:
            add_text(slide, "→", x + 1.22, 2.56, 0.30, 0.32, font_size=16, color_hex=PURPLE, bold=True, align=PP_ALIGN.CENTER)

    bottom_cards = [
        ("통합 워크스페이스", "대시보드, 주간 플래너, 집중 모드를 하나의 세션으로 연결"),
        ("제안형 AI", "AI가 멋대로 바꾸지 않고 구조화된 제안과 재배치 보조에 집중"),
        ("운영 데이터 구조", "일정 블록, 목표, 태스크, 포커스, sync 상태를 함께 다루는 설계"),
    ]
    positions = [0.82, 4.48, 8.14]
    for (title, body), x in zip(bottom_cards, positions):
        add_card(slide, x, 4.70, 3.32, 1.72, fill_hex=WHITE, line_hex=LINE)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, 4.70, 0.10, 1.72, fill_hex=PURPLE, line_hex=None)
        add_text(slide, title, x + 0.22, 4.94, 2.72, 0.28, font_size=13, color_hex=NAVY, bold=True, margin_left=0)
        add_text(slide, body, x + 0.22, 5.28, 2.78, 0.60, font_size=10, color_hex=SLATE, margin_left=0)
    add_footer(slide, "02")


def build_progress_slide(slide) -> None:
    decorate_light(slide)
    add_chip(slide, "CURRENT PROGRESS", 0.82, 0.52, 1.65, 0.34, fill_hex=PURPLE_SOFT, text_hex=PURPLE, font_size=9)
    add_text(
        slide,
        "핵심 사용자 흐름은 화면 단위로 시연 가능한 수준까지 구현했습니다.",
        0.82,
        0.96,
        10.20,
        0.88,
        font_size=21,
        color_hex=NAVY,
        bold=True,
        margin_left=0,
    )
    add_card(slide, 0.82, 1.84, 3.30, 4.86, fill_hex=WHITE, line_hex=LINE)
    add_chip(slide, "진행 상황", 1.04, 2.04, 0.92, 0.30, fill_hex=PURPLE, text_hex=WHITE, font_size=9)
    add_text(slide, "화면과 핵심 액션은\nMVP 데모 가능한 단계", 1.02, 2.46, 2.70, 0.68, font_size=21, color_hex=NAVY, bold=True, margin_left=0, line_spacing=1.15)

    sections = [
        ("구현 화면", "로그인 / 온보딩 / 대시보드 / 주간 플래너 / 집중 모드"),
        ("연결 액션", "일정 블록 생성·수정, 집중 완료·연장·미루기, AI 제안 적용·거절"),
        ("현재 판단", "내부 데모 가능. 운영 안정화는 보완 필요"),
    ]
    section_y = 3.46
    for title, body in sections:
        add_card(slide, 1.02, section_y, 2.86, 0.84, fill_hex=IVORY, line_hex=LINE)
        add_text(slide, title, 1.18, section_y + 0.14, 0.92, 0.20, font_size=10, color_hex=PURPLE, bold=True, margin_left=0)
        add_text(slide, body, 1.18, section_y + 0.36, 2.42, 0.32, font_size=9, color_hex=SLATE, margin_left=0)
        section_y += 0.98

    add_card(slide, 1.02, 6.06, 1.24, 0.44, fill_hex=MINT_SOFT, line_hex="BBDDD3")
    add_card(slide, 2.40, 6.06, 1.24, 0.44, fill_hex=PURPLE_SOFT, line_hex=PURPLE_LINE)
    add_text(slide, "7개 사용자 경로", 1.10, 6.14, 1.05, 0.16, font_size=8, color_hex=SLATE, align=PP_ALIGN.CENTER)
    add_text(slide, "5개 주요 캡처", 2.48, 6.14, 1.05, 0.16, font_size=8, color_hex=SLATE, align=PP_ALIGN.CENTER)

    add_picture_card(slide, SCREENSHOT_DIR / "02-onboarding.png", 4.42, 1.84, 3.95, 2.00, label="온보딩")
    add_picture_card(slide, SCREENSHOT_DIR / "03-dashboard.png", 8.58, 1.84, 3.95, 2.00, label="대시보드")
    add_picture_card(slide, SCREENSHOT_DIR / "04-schedule.png", 4.42, 4.06, 4.30, 2.58, label="주간 플래너")
    add_picture_card(slide, SCREENSHOT_DIR / "05-focus.png", 8.92, 4.06, 3.61, 2.58, label="집중 모드")
    add_footer(slide, "03")


def build_status_slide(slide) -> None:
    decorate_light(slide)
    add_chip(slide, "STATUS CHECK", 0.82, 0.52, 1.22, 0.34, fill_hex=PURPLE_SOFT, text_hex=PURPLE, font_size=9)
    add_text(
        slide,
        "프론트-백 연동은 MVP 수준이며, 외부 연동과 운영 안정화는 아직 남았습니다.",
        0.82,
        0.96,
        8.9,
        0.82,
        font_size=21,
        color_hex=NAVY,
        bold=True,
        margin_left=0,
    )
    add_text(
        slide,
        "현재 저장소에는 넓은 범위의 API와 화면 골격이 있으며, 핵심 차이는 '구현 존재'와 '운영 완성도'입니다.",
        0.82,
        1.76,
        9.4,
        0.25,
        font_size=11,
        color_hex=SLATE,
        margin_left=0,
    )

    add_card(slide, 10.10, 0.84, 1.18, 0.88, fill_hex=WHITE, line_hex=LINE)
    add_card(slide, 11.46, 0.84, 1.18, 0.88, fill_hex=WHITE, line_hex=LINE)
    add_text(slide, "7", 10.10, 0.92, 1.18, 0.28, font_size=24, color_hex=PURPLE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "front routes", 10.10, 1.28, 1.18, 0.14, font_size=8, color_hex=SLATE, align=PP_ALIGN.CENTER)
    add_text(slide, "15", 11.46, 0.92, 1.18, 0.28, font_size=24, color_hex=MINT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "API controllers", 11.46, 1.28, 1.18, 0.14, font_size=8, color_hex=SLATE, align=PP_ALIGN.CENTER)

    columns = [
        (
            "완료",
            MINT,
            MINT_SOFT,
            [
                "주간 플래너 UI",
                "집중 모드 UI",
                "대시보드 핵심 화면",
                "일정 블록 CRUD",
                "포커스 완료 / 연장 / 미루기",
            ],
        ),
        (
            "MVP 연결",
            PURPLE,
            PURPLE_SOFT,
            [
                "로그인 / auth callback",
                "온보딩 bootstrap / answers",
                "AI 재배치 suggestion 흐름",
                "Google sync 상태 표시",
                "Settings API 기본 저장",
            ],
        ),
        (
            "남은 과제",
            AMBER,
            AMBER_SOFT,
            [
                "온보딩 완료 영속화",
                "설정 화면 및 패턴 저장 확장",
                "대시보드 aggregate API",
                "sync 실동작과 scaffold 구분",
                "LLM 품질·설명 강화",
            ],
        ),
    ]

    x_positions = [0.82, 4.46, 8.10]
    for (title, accent, fill_hex, items), x in zip(columns, x_positions):
        add_card(slide, x, 1.94, 3.12, 4.92, fill_hex=WHITE, line_hex=LINE)
        add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, 1.94, 3.12, 0.20, fill_hex=accent, line_hex=None)
        add_text(slide, title, x + 0.20, 2.24, 1.10, 0.28, font_size=16, color_hex=NAVY, bold=True, margin_left=0)
        row_y = 2.76
        for item in items:
            add_card(slide, x + 0.18, row_y, 2.76, 0.58, fill_hex=fill_hex, line_hex=None)
            add_text(slide, item, x + 0.30, row_y + 0.18, 2.44, 0.18, font_size=10, color_hex=SLATE, margin_left=0)
            row_y += 0.74

    add_card(slide, 0.82, 6.78, 11.40, 0.28, fill_hex=PURPLE_SOFT, line_hex=None)
    add_text(
        slide,
        "메모: 대시보드 집계는 프론트 조합 방식이고, sync와 AI 재배치 흐름은 아직 scaffold / MVP 성격이 강합니다.",
        1.00,
        6.82,
        11.00,
        0.18,
        font_size=8,
        color_hex=SLATE,
        margin_left=0,
    )
    add_footer(slide, "04")


def build_roadmap_slide(slide) -> None:
    decorate_dark(slide)
    add_chip(slide, "NEXT MILESTONE", 0.82, 0.56, 1.38, 0.34, fill_hex=PURPLE, text_hex=WHITE, font_size=9)
    add_text(
        slide,
        "다음 보고 시점 목표는 '운영 가능한 연결형 MVP'입니다.",
        0.82,
        1.02,
        8.4,
        0.84,
        font_size=22,
        color_hex=WHITE,
        bold=True,
        margin_left=0,
    )
    add_text(
        slide,
        "남은 일은 분명합니다. 계정 연결과 사용자 상태를 안정화하고, 외부 sync와 AI 제안의 신뢰도를 끌어올려야 합니다.",
        0.82,
        1.86,
        9.5,
        0.28,
        font_size=11,
        color_hex=PURPLE_SOFT,
        margin_left=0,
    )

    phases = [
        ("1. 안정화", PURPLE, "운영형 로그인 / OAuth\n온보딩 완료 저장\n설정 화면 추가"),
        ("2. 연동 완성", MINT, "대시보드 aggregate API\nsync 상태 명확화\nschedule block 기반 focus 보강"),
        ("3. AI 고도화", AMBER, "LLM 정규화 품질 검증\n재배치 설명 payload 강화\n사용자 신뢰도 개선"),
    ]
    x_positions = [0.82, 4.48, 8.14]
    for (title, accent, body), x in zip(phases, x_positions):
        add_card(slide, x, 2.16, 3.20, 3.72, fill_hex=NAVY_2, line_hex="30507B", transparency=0.04)
        add_chip(slide, title, x + 0.22, 2.40, 1.08, 0.32, fill_hex=accent, text_hex=WHITE, font_size=9)
        add_text(slide, "핵심 과제", x + 0.22, 2.92, 0.90, 0.20, font_size=10, color_hex=PURPLE_SOFT, bold=True, margin_left=0)
        add_text(slide, body, x + 0.22, 3.22, 2.60, 1.20, font_size=13, color_hex=WHITE, bold=True, margin_left=0, line_spacing=1.3)
        add_text(slide, "이번 단계가 끝나면 다음 단계로 넘어갈 수 있는 조건이 확보됩니다.", x + 0.22, 5.10, 2.60, 0.46, font_size=9, color_hex=PURPLE_SOFT, margin_left=0)

    add_text(slide, "→", 3.94, 3.66, 0.32, 0.32, font_size=18, color_hex=PURPLE_SOFT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "→", 7.60, 3.66, 0.32, 0.32, font_size=18, color_hex=PURPLE_SOFT, bold=True, align=PP_ALIGN.CENTER)

    add_card(slide, 0.82, 6.30, 11.46, 0.72, fill_hex="203C63", line_hex="30507B")
    add_text(
        slide,
        "중간 결론: UI 골격과 핵심 API는 확보했다. 다음 단계는 데이터 영속화, 외부 연동, AI 신뢰도를 메우는 일이다.",
        1.02,
        6.48,
        11.00,
        0.22,
        font_size=12,
        color_hex=WHITE,
        bold=True,
        margin_left=0,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, "05", dark=True)


def build_presentation() -> Presentation:
    prs = Presentation()
    set_slide_size(prs)
    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.title = "Time Table Mid Progress Report"
    prs.core_properties.subject = "PM Update"
    blank = prs.slide_layouts[6]

    build_cover(prs.slides.add_slide(blank))
    build_goal_slide(prs.slides.add_slide(blank))
    build_progress_slide(prs.slides.add_slide(blank))
    build_status_slide(prs.slides.add_slide(blank))
    build_roadmap_slide(prs.slides.add_slide(blank))
    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(OUTPUT_PATH)
    print(f"saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
