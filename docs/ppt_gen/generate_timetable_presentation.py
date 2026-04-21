from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUTPUT_PATH = ROOT / "docs" / "TimeTable_Project_Presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x10, 0x10, 0x10)
TEXT_GRAY = RGBColor(0x4F, 0x4F, 0x4F)
LINE_GRAY = RGBColor(0xD9, 0xD9, 0xD9)
FILL_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
SHADE_GRAY = RGBColor(0xEB, 0xEB, 0xEB)

LEFT = Inches(0.96)
RIGHT = Inches(12.37)
TOP = Inches(0.70)
CONTENT_W = Inches(11.35)

TITLE_SIZE = 24
SUBTITLE_SIZE = 18
BODY_SIZE = 16
CAPTION_SIZE = 14
PAGE_SIZE = 11
FONT = "Malgun Gothic"


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=BODY_SIZE,
    bold=False,
    color=BLACK,
    align=PP_ALIGN.LEFT,
    margin=0.02,
    line_spacing=1.25,
    space_after=2,
    italic=False,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_paragraphs(
    slide,
    lines,
    x,
    y,
    w,
    h,
    *,
    size=BODY_SIZE,
    color=BLACK,
    bold_first=False,
    line_spacing=1.25,
    space_after=4,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    first = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold_first and idx == 0
    return box


def add_box(
    slide,
    x,
    y,
    w,
    h,
    *,
    fill=WHITE,
    line=LINE_GRAY,
    line_width=0.8,
):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    return shape


def add_rule(slide, x, y, w, h=Inches(0.02), *, color=LINE_GRAY):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_page(slide, number: int) -> None:
    add_text(
        slide,
        str(number),
        Inches(12.55),
        Inches(6.92),
        Inches(0.28),
        Inches(0.18),
        size=PAGE_SIZE,
        color=TEXT_GRAY,
        align=PP_ALIGN.RIGHT,
        margin=0,
    )


def header(slide, title: str, subtitle: str, page: int) -> None:
    set_bg(slide)
    add_text(slide, title, LEFT, TOP, Inches(8.9), Inches(0.38), size=TITLE_SIZE, bold=True)
    add_text(
        slide,
        subtitle,
        LEFT,
        Inches(1.40),
        Inches(9.9),
        Inches(0.38),
        size=BODY_SIZE,
        color=TEXT_GRAY,
    )
    add_page(slide, page)


def add_cover_motif(slide) -> None:
    x = Inches(10.20)
    y = Inches(5.05)
    cell_w = Inches(0.55)
    cell_h = Inches(0.32)
    for row in range(4):
        for col in range(4):
            fill = WHITE
            if (row, col) in {(1, 1), (2, 2), (0, 3)}:
                fill = SHADE_GRAY
            add_box(
                slide,
                x + cell_w * col,
                y + cell_h * row,
                cell_w - Inches(0.02),
                cell_h - Inches(0.02),
                fill=fill,
                line=LINE_GRAY,
                line_width=0.6,
            )


def cover_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(
        slide,
        "인자부터 니 이름은 춘식이여  ·  이근산  ·  2026.04.16",
        LEFT,
        Inches(0.82),
        Inches(6.2),
        Inches(0.22),
        size=CAPTION_SIZE,
        color=TEXT_GRAY,
    )
    add_text(slide, "타임테이블", LEFT, Inches(2.50), Inches(4.2), Inches(0.36), size=TITLE_SIZE, bold=True)
    add_text(
        slide,
        "AI 기반 일정 관리 비서 서비스",
        LEFT,
        Inches(2.95),
        Inches(4.6),
        Inches(0.28),
        size=SUBTITLE_SIZE,
        color=TEXT_GRAY,
    )
    add_cover_motif(slide)
    add_page(slide, 1)


def slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "일정 관리의 진짜 문제", "일정 관리의 본질적인 어려움은 기록이 아니라 조율이다.", 2)
    add_paragraphs(
        slide,
        [
            "계획을 세우는 것보다 바뀐 일정을 다시 맞추는 과정이 더 힘들다.",
            "회의, 과제, 개인 일정이 겹칠수록 직접 재계산해야 하는 부담이 커진다.",
            "사용자는 저장보다 재조정에서 더 큰 피로를 느낀다.",
        ],
        LEFT,
        Inches(2.05),
        Inches(6.2),
        Inches(2.4),
        size=BODY_SIZE,
        color=BLACK,
        line_spacing=1.28,
        space_after=8,
    )
    dx = Inches(8.25)
    dy = Inches(2.15)
    add_box(slide, dx, dy, Inches(1.55), Inches(0.62), fill=FILL_GRAY)
    add_box(slide, dx + Inches(1.75), dy + Inches(0.28), Inches(1.55), Inches(0.62), fill=WHITE)
    add_box(slide, dx + Inches(0.82), dy + Inches(1.10), Inches(1.55), Inches(0.62), fill=WHITE)
    add_text(slide, "캘린더", dx + Inches(0.20), dy + Inches(0.18), Inches(1.1), Inches(0.2), size=SUBTITLE_SIZE, bold=True)
    add_text(slide, "태스크", dx + Inches(1.95), dy + Inches(0.46), Inches(1.1), Inches(0.2), size=SUBTITLE_SIZE, bold=True)
    add_text(slide, "메모", dx + Inches(1.02), dy + Inches(1.28), Inches(0.9), Inches(0.2), size=SUBTITLE_SIZE, bold=True)
    add_text(slide, "↘", dx + Inches(1.08), dy + Inches(0.68), Inches(0.3), Inches(0.2), size=BODY_SIZE, color=TEXT_GRAY)
    add_text(slide, "↙", dx + Inches(2.06), dy + Inches(0.86), Inches(0.3), Inches(0.2), size=BODY_SIZE, color=TEXT_GRAY)
    add_text(slide, "사용자가 직접 다시 계산", dx - Inches(0.05), dy + Inches(2.00), Inches(3.6), Inches(0.24), size=BODY_SIZE, color=TEXT_GRAY)


def add_step_flow(slide, labels, x, y, total_w, box_h, *, text_size=BODY_SIZE, arrow_gap=0.16):
    count = len(labels)
    box_w = (total_w - Inches(arrow_gap) * (count - 1)) / count
    for idx, label in enumerate(labels):
        bx = x + idx * (box_w + Inches(arrow_gap))
        add_box(slide, bx, y, box_w, box_h, fill=WHITE)
        add_text(slide, label, bx + Inches(0.12), y + Inches(0.15), box_w - Inches(0.24), box_h - Inches(0.20), size=text_size, align=PP_ALIGN.CENTER)
        if idx < count - 1:
            add_text(
                slide,
                "→",
                bx + box_w + Inches(0.02),
                y + Inches(0.12),
                Inches(0.18),
                Inches(0.20),
                size=BODY_SIZE,
                color=TEXT_GRAY,
                align=PP_ALIGN.CENTER,
                margin=0,
            )


def slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "사용자가 실제로 겪는 불편", "기존 흐름은 변경이 생기는 순간 다시 시작된다.", 3)
    add_step_flow(
        slide,
        ["캘린더 저장", "태스크 정리", "변경 발생", "재구성", "우선순위 재판단"],
        LEFT,
        Inches(2.65),
        Inches(10.9),
        Inches(0.84),
        text_size=CAPTION_SIZE,
    )
    add_rule(slide, LEFT, Inches(4.55), CONTENT_W)
    add_text(
        slide,
        "일정 관리가 번거로운 이유는 일정이 많아서가 아니라, 계속 다시 짜야 해서다.",
        LEFT,
        Inches(4.85),
        CONTENT_W,
        Inches(0.30),
        size=SUBTITLE_SIZE,
        bold=True,
    )


def draw_day_schedule(slide, x, y, title, entries):
    add_text(slide, title, x, y, Inches(1.8), Inches(0.24), size=SUBTITLE_SIZE, bold=True)
    add_box(slide, x, y + Inches(0.38), Inches(2.95), Inches(3.63), fill=WHITE)
    row_h = Inches(0.54)
    time_w = Inches(0.68)
    for idx, (time_label, block_label, shaded) in enumerate(entries):
        ry = y + Inches(0.38) + row_h * idx
        add_rule(slide, x, ry, Inches(2.95), Inches(0.01))
        add_rule(slide, x + time_w, ry, Inches(0.01), row_h)
        if shaded:
            add_box(
                slide,
                x + time_w,
                ry + Inches(0.01),
                Inches(2.27),
                row_h - Inches(0.02),
                fill=SHADE_GRAY,
                line=SHADE_GRAY,
                line_width=0.1,
            )
        add_text(slide, time_label, x + Inches(0.06), ry + Inches(0.14), Inches(0.50), Inches(0.18), size=CAPTION_SIZE, color=TEXT_GRAY)
        add_text(slide, block_label, x + time_w + Inches(0.10), ry + Inches(0.12), Inches(2.05), Inches(0.24), size=CAPTION_SIZE)
    add_rule(slide, x, y + Inches(0.38) + row_h * len(entries), Inches(2.95), Inches(0.01))


def slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "문제 상황 예시", "작은 변경 3개가 하루 전체를 다시 계산하게 만든다.", 4)
    add_text(
        slide,
        "오후 3시 팀플 회의가 1시간 밀림   ·   저녁 운동 일정이 추가됨   ·   과제 마감이 가까워짐",
        LEFT,
        Inches(1.98),
        CONTENT_W,
        Inches(0.24),
        size=CAPTION_SIZE,
        color=TEXT_GRAY,
    )
    before_entries = [
        ("09", "수업", False),
        ("11", "과제 진행", False),
        ("13", "점심 / 이동", False),
        ("15", "팀플 회의", False),
        ("17", "자유 시간", False),
        ("19", "과제 마무리", False),
    ]
    after_entries = [
        ("09", "수업", False),
        ("11", "과제 진행", False),
        ("13", "점심 / 이동", False),
        ("15", "팀플 회의 1시간 지연", True),
        ("17", "운동 추가", True),
        ("19", "과제 우선 재배치", True),
    ]
    draw_day_schedule(slide, Inches(1.45), Inches(2.45), "변경 전", before_entries)
    draw_day_schedule(slide, Inches(8.00), Inches(2.45), "변경 후", after_entries)


def slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "기존 도구의 한계", "각 도구는 일부만 해결하고, 전체 일정은 다시 사람이 맞춘다.", 5)
    boxes = [
        (LEFT, Inches(2.15), "캘린더", "일정 저장은 쉽지만\n우선순위 판단은 약함"),
        (Inches(7.10), Inches(2.15), "태스크 앱", "할 일 정리는 되지만\n시간표와 바로 연결되지 않음"),
        (LEFT, Inches(4.40), "메모 앱", "정보는 적기 쉽지만\n배치 기준이 없음"),
        (Inches(7.10), Inches(4.40), "범용 AI", "답변은 가능하지만\n개인 일정 맥락이 이어지지 않음"),
    ]
    for x, y, title, body in boxes:
        add_box(slide, x, y, Inches(4.45), Inches(1.55), fill=FILL_GRAY)
        add_text(slide, title, x + Inches(0.18), y + Inches(0.18), Inches(2.4), Inches(0.22), size=SUBTITLE_SIZE, bold=True)
        add_text(slide, body, x + Inches(0.18), y + Inches(0.56), Inches(3.9), Inches(0.50), size=CAPTION_SIZE, color=TEXT_GRAY)
    add_text(
        slide,
        "기존 도구는 저장에는 강하지만 조율에는 약하다.",
        LEFT,
        Inches(6.35),
        Inches(6.8),
        Inches(0.28),
        size=SUBTITLE_SIZE,
        bold=True,
    )


def slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "왜 기존 방식이 번거로운가", "정보가 흩어져 있고, 실행 판단은 자동화되지 않는다.", 6)
    y_positions = [Inches(2.20), Inches(3.50), Inches(4.80)]
    items = [
        ("1", "일정과 할 일이 분리되어 있다."),
        ("2", "우선순위와 시간표가 연결되지 않는다."),
        ("3", "실제 실행 가능한 수준까지 자동화되지 않는다."),
    ]
    for (num, text), y in zip(items, y_positions):
        add_text(slide, num, LEFT, y, Inches(0.30), Inches(0.20), size=CAPTION_SIZE, bold=True, color=TEXT_GRAY)
        add_text(slide, text, LEFT + Inches(0.42), y - Inches(0.02), Inches(8.8), Inches(0.26), size=SUBTITLE_SIZE)
        add_rule(slide, LEFT, y + Inches(0.34), Inches(9.5))
    add_text(
        slide,
        "기록은 앱이 하지만, 판단은 여전히 사용자가 한다.",
        LEFT,
        Inches(6.18),
        Inches(8.3),
        Inches(0.28),
        size=SUBTITLE_SIZE,
        bold=True,
    )


def slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "시장의 빈칸", "시장은 기록과 추천은 많지만, 일정 재배치는 아직 비어 있다.", 7)
    stages = [
        ("기록형 도구", False),
        ("추천형 AI", False),
        ("조율형 비서", True),
    ]
    x = LEFT
    gap = Inches(0.40)
    width = Inches(3.45)
    for title, shaded in stages:
        add_box(slide, x, Inches(2.65), width, Inches(0.78), fill=SHADE_GRAY if shaded else WHITE)
        add_text(slide, title, x + Inches(0.18), Inches(2.92), width - Inches(0.36), Inches(0.24), size=SUBTITLE_SIZE, bold=True, align=PP_ALIGN.CENTER)
        x += width + gap
    add_text(slide, "→", Inches(4.70), Inches(2.92), Inches(0.20), Inches(0.18), size=BODY_SIZE, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, "→", Inches(8.57), Inches(2.92), Inches(0.20), Inches(0.18), size=BODY_SIZE, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    add_paragraphs(
        slide,
        [
            "일정 저장 도구는 많다.",
            "AI 추천 기능도 늘고 있다.",
            "그러나 일정 맥락을 이해하고 재배치하는 서비스는 드물다.",
        ],
        LEFT,
        Inches(4.45),
        Inches(10.4),
        Inches(1.1),
        size=BODY_SIZE,
        color=BLACK,
        line_spacing=1.28,
        space_after=8,
    )


def slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "타임테이블이 해결하는 방식", "입력은 자연어로 받고, 출력은 실행 가능한 시간표로 제공한다.", 8)
    add_paragraphs(
        slide,
        [
            "구글 캘린더와 구글 태스크를 연동한다.",
            "자연어 입력으로 일정을 받는다.",
            "AI와 룰베이스 조합으로 일정을 해석하고 조정한다.",
            "자동 시간표 생성과 변경 반영을 지원한다.",
        ],
        LEFT,
        Inches(2.05),
        Inches(6.15),
        Inches(2.8),
        size=BODY_SIZE,
        color=BLACK,
        line_spacing=1.28,
        space_after=8,
    )
    x = Inches(8.25)
    labels = ["캘린더 / 태스크", "자연어 입력", "AI + 룰 해석", "시간표 생성\n변경 반영"]
    for idx, label in enumerate(labels):
        by = Inches(2.18) + Inches(0.95) * idx
        add_box(slide, x, by, Inches(3.15), Inches(0.62), fill=FILL_GRAY if idx == 2 else WHITE)
        add_text(slide, label, x + Inches(0.18), by + Inches(0.14), Inches(2.75), Inches(0.28), size=CAPTION_SIZE, align=PP_ALIGN.CENTER)
        if idx < len(labels) - 1:
            add_text(slide, "↓", x + Inches(1.45), by + Inches(0.66), Inches(0.24), Inches(0.18), size=BODY_SIZE, color=TEXT_GRAY, align=PP_ALIGN.CENTER)


def slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "서비스 목표", "사용자는 조건을 알려주고, 시스템은 배치와 재조정을 맡는다.", 9)
    add_text(slide, "사용자가 알려주는 것", LEFT, Inches(2.20), Inches(2.8), Inches(0.24), size=SUBTITLE_SIZE, bold=True)
    add_text(slide, "시스템이 대신하는 것", Inches(7.15), Inches(2.20), Inches(2.8), Inches(0.24), size=SUBTITLE_SIZE, bold=True)
    add_rule(slide, Inches(6.52), Inches(2.08), Inches(0.02), Inches(3.25))
    add_paragraphs(slide, ["일정", "마감", "중요도", "조건"], LEFT, Inches(2.75), Inches(2.2), Inches(1.8), size=BODY_SIZE, color=BLACK, line_spacing=1.28, space_after=8)
    add_paragraphs(slide, ["배치", "재조정", "추천", "현재 할 일 제시"], Inches(7.15), Inches(2.75), Inches(3.2), Inches(1.8), size=BODY_SIZE, color=BLACK, line_spacing=1.28, space_after=8)
    add_text(slide, "타임테이블은 계획 도구가 아니라 계획을 위임하는 서비스다.", LEFT, Inches(6.10), Inches(8.8), Inches(0.28), size=SUBTITLE_SIZE, bold=True)


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "핵심 제안 가치", "서비스 가치는 입력·판단·재조정·실행 전환 비용을 함께 줄이는 데 있다.", 10)
    items = [
        (LEFT, Inches(2.20), "일정 생성 자동화", "초기 계획 부담 감소"),
        (Inches(7.05), Inches(2.20), "변경 반영 자동화", "변경 대응 부담 감소"),
        (LEFT, Inches(4.55), "우선순위 판단 보조", "판단 속도 개선"),
        (Inches(7.05), Inches(4.55), "실행 전환 비용 감소", "지금 할 일 전환 지원"),
    ]
    for x, y, title, body in items:
        add_box(slide, x, y, Inches(4.35), Inches(1.55), fill=WHITE)
        add_text(slide, title, x + Inches(0.18), y + Inches(0.24), Inches(3.8), Inches(0.24), size=SUBTITLE_SIZE, bold=True)
        add_text(slide, body, x + Inches(0.18), y + Inches(0.70), Inches(3.8), Inches(0.20), size=CAPTION_SIZE, color=TEXT_GRAY)
    add_text(slide, "핵심은 예쁜 캘린더가 아니라, 오늘 실제로 움직일 수 있는 일정표다.", LEFT, Inches(6.33), Inches(10.7), Inches(0.28), size=SUBTITLE_SIZE, bold=True)


def slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "누가 이 서비스를 필요로 하는가", "핵심 타깃은 일정 수보다 조율 부담이 큰 사용자다.", 11)
    rows = [
        "일정이 많은 사람",
        "일정 변경이 잦은 사람",
        "직접 조율하는 과정이 피곤한 사람",
        "해야 할 일을 즉시 알고 싶은 사람",
    ]
    y = Inches(2.18)
    for idx, row in enumerate(rows, start=1):
        add_text(slide, f"{idx:02d}", LEFT, y, Inches(0.34), Inches(0.18), size=CAPTION_SIZE, bold=True, color=TEXT_GRAY)
        add_text(slide, row, LEFT + Inches(0.55), y - Inches(0.02), Inches(9.2), Inches(0.24), size=SUBTITLE_SIZE)
        add_rule(slide, LEFT, y + Inches(0.38), Inches(9.8))
        y += Inches(0.95)
    add_text(slide, "핵심 타겟은 일정은 많은데 조율까지 직접 하기 싫은 사람들이다.", LEFT, Inches(6.20), Inches(10.4), Inches(0.28), size=SUBTITLE_SIZE, bold=True)


def slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "핵심 사용자 유형", "학생과 직장인 모두 일정 수보다 조율 부담이 큰 순간에 이 서비스를 필요로 한다.", 12)
    add_box(slide, LEFT, Inches(2.15), Inches(4.9), Inches(2.95), fill=FILL_GRAY)
    add_box(slide, Inches(7.10), Inches(2.15), Inches(4.9), Inches(2.95), fill=FILL_GRAY)
    add_text(slide, "학생", LEFT + Inches(0.18), Inches(2.38), Inches(1.2), Inches(0.22), size=SUBTITLE_SIZE, bold=True)
    add_text(slide, "직장인 / 멀티 프로젝트\n사용자", Inches(7.28), Inches(2.36), Inches(3.4), Inches(0.50), size=16, bold=True)
    add_paragraphs(slide, ["수업, 과제, 시험, 개인 일정이 동시에 존재", "마감과 공부 시간을 함께 관리해야 함"], LEFT + Inches(0.18), Inches(2.84), Inches(4.2), Inches(1.3), size=15, color=BLACK, line_spacing=1.22, space_after=8)
    add_paragraphs(slide, ["회의, 업무, 개인 일정이 수시로 변경됨", "집중 시간 확보와 우선순위 재조정이 중요함"], Inches(7.28), Inches(3.12), Inches(4.2), Inches(1.3), size=15, color=BLACK, line_spacing=1.22, space_after=8)


def slide_13(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "사용자 니즈 정리", "사용자가 원하는 것은 입력 도구가 아니라 실행 보조다.", 13)
    needs = [
        "입력은 간단해야 한다.",
        "일정 충돌은 자동으로 해결되어야 한다.",
        "지금 해야 할 일을 즉시 알려줘야 한다.",
        "변경이 생기면 전체 흐름이 다시 정리되어야 한다.",
    ]
    y = Inches(2.18)
    for need in needs:
        add_text(slide, "•", LEFT, y, Inches(0.22), Inches(0.18), size=BODY_SIZE, bold=True)
        add_text(slide, need, LEFT + Inches(0.32), y - Inches(0.02), Inches(9.6), Inches(0.24), size=SUBTITLE_SIZE)
        y += Inches(0.92)
    add_text(slide, "사용자가 원하는 것은 입력 도구가 아니라 실행 보조다.", LEFT, Inches(6.10), Inches(8.5), Inches(0.28), size=SUBTITLE_SIZE, bold=True)


def slide_14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "서비스 작동 구조", "입력부터 재배치와 현재 추천까지 하나의 흐름으로 연결된다.", 14)
    labels = [
        "일정 및\n할 일 입력",
        "외부 일정\n데이터 불러오기",
        "우선순위와\n제약 조건 해석",
        "시간표\n자동 생성",
        "변경 발생 시\n재배치",
        "현재 시점\n추천 제공",
    ]
    add_step_flow(slide, labels, LEFT, Inches(2.70), CONTENT_W, Inches(0.86), text_size=CAPTION_SIZE)


def slide_15(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "입력 데이터와 판단 요소", "시간표를 만들기 위해서는 일정 정보와 배치 조건을 함께 봐야 한다.", 15)
    add_text(slide, "입력 데이터", LEFT, Inches(2.12), Inches(2.2), Inches(0.22), size=SUBTITLE_SIZE, bold=True)
    add_text(slide, "판단 요소", Inches(7.40), Inches(2.12), Inches(2.2), Inches(0.22), size=SUBTITLE_SIZE, bold=True)
    add_rule(slide, Inches(6.55), Inches(2.02), Inches(0.02), Inches(3.90))
    add_paragraphs(slide, ["일정명", "마감 시간", "예상 소요시간", "중요도", "반복 여부", "선호 시간대", "캘린더 / 태스크 데이터"], LEFT, Inches(2.60), Inches(4.7), Inches(3.4), size=BODY_SIZE, color=BLACK, line_spacing=1.24, space_after=6)
    add_paragraphs(slide, ["충돌 여부", "버퍼 시간", "집중 가능 시간", "일정 분할 가능성"], Inches(7.40), Inches(2.60), Inches(3.7), Inches(2.4), size=BODY_SIZE, color=BLACK, line_spacing=1.28, space_after=8)


def slide_16(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "AI와 룰베이스의 역할 분담", "AI는 해석과 설명을 맡고, 룰베이스는 배치와 충돌 계산을 맡는다.", 16)
    add_text(slide, "AI", LEFT, Inches(2.10), Inches(1.2), Inches(0.22), size=SUBTITLE_SIZE, bold=True)
    add_text(slide, "룰베이스", Inches(7.05), Inches(2.10), Inches(1.8), Inches(0.22), size=SUBTITLE_SIZE, bold=True)
    add_rule(slide, Inches(6.42), Inches(2.02), Inches(0.02), Inches(3.65))
    add_paragraphs(slide, ["자연어 이해", "변경 의도 해석", "대화형 조정", "추천 설명 생성"], LEFT, Inches(2.60), Inches(3.8), Inches(2.6), size=BODY_SIZE, color=BLACK, line_spacing=1.28, space_after=8)
    add_paragraphs(slide, ["충돌 판정", "우선순위 반영", "시간 슬롯 배치", "일정 제약 처리"], Inches(7.05), Inches(2.60), Inches(3.4), Inches(2.6), size=BODY_SIZE, color=BLACK, line_spacing=1.28, space_after=8)
    add_text(slide, "AI만으로 판단하지 않고 규칙 로직과 결합해 신뢰성을 확보한다.", LEFT, Inches(6.12), Inches(9.0), Inches(0.28), size=SUBTITLE_SIZE, bold=True)


def slide_17(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "사용자 경험의 핵심 기능 (1)", "입력과 계획 생성은 자연어와 자동 배치 경험으로 연결된다.", 17)
    add_text(slide, "1) 자연어 기반 일정 입력", LEFT, Inches(2.10), Inches(4.8), Inches(0.24), size=SUBTITLE_SIZE, bold=True)
    add_paragraphs(slide, ['“내일 오후 3시에 회의 넣어줘”', '“이번 주 안에 과제 3시간 확보해줘”'], LEFT, Inches(2.58), Inches(4.8), Inches(0.84), size=15, color=TEXT_GRAY, line_spacing=1.22, space_after=6)
    add_rule(slide, LEFT, Inches(4.05), CONTENT_W)
    add_text(slide, "2) AI 자동 시간표 생성", LEFT, Inches(4.40), Inches(4.8), Inches(0.24), size=SUBTITLE_SIZE, bold=True)
    add_paragraphs(slide, ["우선순위, 마감, 소요시간 반영", "주간 및 일간 계획 자동 배치"], LEFT, Inches(4.88), Inches(5.4), Inches(0.78), size=BODY_SIZE, color=BLACK, line_spacing=1.25, space_after=8)


def slide_18(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "사용자 경험의 핵심 기능 (2)", "변경 대응, 현재 추천, 외부 연동이 실제 사용 경험을 완성한다.", 18)
    sections = [
        ("3) 변경 반영 및 자동 재조정", "지연, 취소, 추가 일정을 반영하고 전체 흐름을 다시 배치"),
        ("4) 지금 해야 할 일 제시", "현재 시점 기준 최적 작업을 추천하고 이유를 함께 설명"),
        ("5) 캘린더 / 태스크 연동", "기존 도구를 버리지 않고 그대로 활용"),
    ]
    y = Inches(2.12)
    for idx, (title, body) in enumerate(sections):
        add_text(slide, title, LEFT, y, Inches(5.5), Inches(0.24), size=SUBTITLE_SIZE, bold=True)
        add_text(slide, body, LEFT, y + Inches(0.42), Inches(10.1), Inches(0.28), size=BODY_SIZE, color=TEXT_GRAY)
        if idx < len(sections) - 1:
            add_rule(slide, LEFT, y + Inches(0.98), CONTENT_W)
        y += Inches(1.35)


def slide_19(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "기능 요약: 사용자 가치", "기능은 결국 입력·계획·변경·실행의 부담을 줄이는 방향으로 연결된다.", 19)
    rows = [
        ("자연어 입력", "입력 부담 감소"),
        ("자동 생성", "계획 수립 부담 감소"),
        ("자동 재조정", "변경 대응 부담 감소"),
        ("지금 할 일 안내", "실행 부담 감소"),
    ]
    add_text(slide, "기능", LEFT, Inches(2.20), Inches(1.2), Inches(0.20), size=CAPTION_SIZE, bold=True, color=TEXT_GRAY)
    add_text(slide, "줄어드는 부담", Inches(7.20), Inches(2.20), Inches(2.2), Inches(0.20), size=CAPTION_SIZE, bold=True, color=TEXT_GRAY)
    add_rule(slide, Inches(6.25), Inches(2.10), Inches(0.02), Inches(3.05))
    y = Inches(2.62)
    for left_text, right_text in rows:
        add_text(slide, left_text, LEFT, y, Inches(4.8), Inches(0.22), size=SUBTITLE_SIZE, bold=True)
        add_text(slide, right_text, Inches(7.20), y, Inches(3.8), Inches(0.22), size=SUBTITLE_SIZE)
        y += Inches(0.72)
    add_text(slide, "타임테이블은 입력 부담, 판단 부담, 재조정 부담을 모두 줄인다.", LEFT, Inches(6.05), Inches(9.6), Inches(0.28), size=SUBTITLE_SIZE, bold=True)


def slide_20(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "서비스 이용 시나리오", "처음 쓰는 흐름과 변경 이후의 흐름이 하나의 사용 경험으로 이어진다.", 20)
    add_text(slide, "초기 사용", LEFT, Inches(2.08), Inches(1.2), Inches(0.18), size=CAPTION_SIZE, bold=True, color=TEXT_GRAY)
    add_step_flow(
        slide,
        ["1. 로그인 및\n구글 연동", "2. 데이터\n불러오기", "3. 일정과\n할 일 입력", "4. 초기 시간표\n생성"],
        LEFT,
        Inches(2.40),
        CONTENT_W,
        Inches(0.82),
        text_size=CAPTION_SIZE,
    )
    add_text(slide, "이후 사용", LEFT, Inches(4.38), Inches(1.2), Inches(0.18), size=CAPTION_SIZE, bold=True, color=TEXT_GRAY)
    add_step_flow(
        slide,
        ["5. 일정표\n확인", "6. 자연어로\n수정 요청", "7. 자동으로\n재배치", "8. 지금 할 일\n제시"],
        LEFT,
        Inches(4.70),
        CONTENT_W,
        Inches(0.82),
        text_size=CAPTION_SIZE,
    )


def slide_21(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "하루 사용 시나리오 예시", "하루 중 변경이 생길 때마다 서비스는 시간표와 현재 행동을 다시 맞춘다.", 21)
    line_x = Inches(2.05)
    add_rule(slide, line_x, Inches(2.18), Inches(0.02), Inches(3.7))
    items = [
        ("오전", "“오늘 할 일 정리해줘”"),
        ("오후", "“회의가 1시간 밀렸어”"),
        ("저녁", "“운동 대신 과제 먼저 할래”"),
        ("결과", "AI가 일정을 다시 조정하고 지금 해야 할 일을 안내"),
    ]
    y = Inches(2.15)
    for label, body in items:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, line_x - Inches(0.08), y + Inches(0.08), Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLACK
        dot.line.fill.background()
        add_text(slide, label, Inches(2.45), y, Inches(0.8), Inches(0.22), size=SUBTITLE_SIZE, bold=True)
        add_text(slide, body, Inches(3.55), y, Inches(6.8), Inches(0.36), size=BODY_SIZE, color=TEXT_GRAY)
        y += Inches(1.02)


def slide_22(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "화면 설계 개요", "화면은 많아 보이는 것보다 역할이 명확해야 한다.", 22)
    rows = [
        ("로그인", "계정 연결과 서비스 진입"),
        ("주간 일정표", "전체 흐름과 주간 배치 확인"),
        ("오늘 일정표", "현재 실행 계획 확인"),
        ("지금 무엇을 해야 하는지", "현재 추천 작업과 추천 이유 표시"),
        ("변경 입력창", "자연어 수정 입력과 결과 반영"),
        ("설정 / 관리", "선호 시간대와 동기화 설정"),
    ]
    add_text(slide, "화면", LEFT, Inches(2.12), Inches(1.1), Inches(0.18), size=CAPTION_SIZE, bold=True, color=TEXT_GRAY)
    add_text(slide, "역할", Inches(6.00), Inches(2.12), Inches(1.1), Inches(0.18), size=CAPTION_SIZE, bold=True, color=TEXT_GRAY)
    add_rule(slide, Inches(5.55), Inches(2.02), Inches(0.02), Inches(3.85))
    y = Inches(2.52)
    for screen, role in rows:
        add_text(slide, screen, LEFT, y, Inches(4.0), Inches(0.20), size=BODY_SIZE, bold=True)
        add_text(slide, role, Inches(6.00), y, Inches(5.3), Inches(0.22), size=BODY_SIZE, color=TEXT_GRAY)
        y += Inches(0.62)


def slide_23(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "핵심 화면 구성", "주간 흐름, 현재 추천, 변경 입력, 설정 관리를 한눈에 구분되게 설계한다.", 23)
    items = [
        (LEFT, Inches(2.10), "1) 주간 / 오늘 일정표", "전체 흐름과 현재 실행 계획을 동시에 확인"),
        (Inches(7.05), Inches(2.10), "2) 지금 무엇을 해야 하는지", "현재 추천 작업 1순위와 추천 이유 표시"),
        (LEFT, Inches(4.42), "3) 변경 입력창", "자연어 수정 입력과 반영 결과 미리보기"),
        (Inches(7.05), Inches(4.42), "4) 설정 / 관리", "선호 집중 시간대, 휴식 간격, 동기화 설정"),
    ]
    for x, y, title, body in items:
        add_box(slide, x, y, Inches(4.45), Inches(1.65), fill=FILL_GRAY)
        add_text(slide, title, x + Inches(0.18), y + Inches(0.18), Inches(3.9), Inches(0.24), size=SUBTITLE_SIZE, bold=True)
        add_text(slide, body, x + Inches(0.18), y + Inches(0.62), Inches(4.0), Inches(0.46), size=CAPTION_SIZE, color=TEXT_GRAY)


def slide_24(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "AI 기반 일정 조율 구조", "상단은 시스템 구성, 하단은 처리 흐름으로 분리해 본다.", 24)
    top_items = [
        (LEFT, Inches(2.08), Inches(1.45), Inches(0.56), "Frontend"),
        (Inches(3.05), Inches(2.08), Inches(1.65), Inches(0.56), "Backend API"),
        (Inches(5.25), Inches(2.08), Inches(1.75), Inches(0.56), "AI Orchestrator"),
        (Inches(7.55), Inches(2.08), Inches(1.55), Inches(0.56), "Rule Engine"),
        (Inches(9.65), Inches(2.08), Inches(1.65), Inches(0.56), "Database"),
        (Inches(6.30), Inches(3.12), Inches(1.95), Inches(0.74), "Google Calendar\nAPI"),
        (Inches(8.55), Inches(3.12), Inches(1.95), Inches(0.74), "Google Tasks\nAPI"),
    ]
    for x, y, w, h, label in top_items:
        add_box(slide, x, y, w, h, fill=WHITE)
        add_text(slide, label, x + Inches(0.08), y + Inches(0.12), w - Inches(0.16), h - Inches(0.12), size=CAPTION_SIZE, align=PP_ALIGN.CENTER)
    add_rule(slide, LEFT, Inches(4.38), CONTENT_W)
    add_step_flow(
        slide,
        ["1. 사용자 입력\n수신", "2. 외부 일정\n동기화", "3. AI가\n의도 해석", "4. 룰엔진이\n배치 계산", "5. 결과 저장\n화면 반영"],
        LEFT,
        Inches(4.78),
        CONTENT_W,
        Inches(0.82),
        text_size=CAPTION_SIZE,
        arrow_gap=0.13,
    )


def slide_25(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "개발 일정", "요구사항 정의부터 MVP 시연까지 단계적으로 진행한다.", 25)
    phases = [
        ("1단계", "요구사항 정의", "사용자 시나리오 정리"),
        ("2단계", "UI / UX 설계", "주요 화면 와이어프레임 제작"),
        ("3단계", "핵심 기능 개발", "연동 및 자동 배치 기능 구현"),
        ("4단계", "테스트 및 보완", "충돌 케이스 검증"),
        ("5단계", "시연 및 배포", "MVP 시연 및 피드백 반영"),
    ]
    width = Inches(2.05)
    gap = Inches(0.22)
    x = LEFT
    for phase, title, detail in phases:
        add_rule(slide, x, Inches(2.18), width)
        add_text(slide, phase, x, Inches(2.42), width, Inches(0.18), size=CAPTION_SIZE, bold=True, color=TEXT_GRAY)
        add_text(slide, title, x, Inches(2.78), width, Inches(0.24), size=SUBTITLE_SIZE, bold=True)
        add_text(slide, detail, x, Inches(3.35), width, Inches(0.50), size=CAPTION_SIZE, color=TEXT_GRAY)
        x += width + gap


def slide_26(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "MVP 범위와 확장 방향", "먼저 핵심 자동 조율 경험을 만들고, 이후 협업과 분석으로 확장한다.", 26)
    add_text(slide, "MVP 범위", LEFT, Inches(2.12), Inches(1.8), Inches(0.22), size=SUBTITLE_SIZE, bold=True)
    add_text(slide, "확장 방향", Inches(7.10), Inches(2.12), Inches(1.8), Inches(0.22), size=SUBTITLE_SIZE, bold=True)
    add_rule(slide, Inches(6.42), Inches(2.02), Inches(0.02), Inches(3.90))
    add_paragraphs(slide, ["구글 캘린더 연동", "구글 태스크 연동", "자연어 일정 입력", "자동 시간표 생성", "변경 반영", "지금 해야 할 일 제시"], LEFT, Inches(2.55), Inches(4.3), Inches(2.8), size=BODY_SIZE, color=BLACK, line_spacing=1.22, space_after=6)
    add_paragraphs(slide, ["팀 일정 공동 조율", "이동시간 반영", "생산성 분석", "음성 비서 인터페이스"], Inches(7.10), Inches(2.55), Inches(4.0), Inches(2.8), size=BODY_SIZE, color=BLACK, line_spacing=1.28, space_after=8)


def slide_27(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "기대효과와 결론", "마지막으로 이 서비스가 줄이는 부담과 남기는 가치를 정리한다.", 27)
    items = [
        (LEFT, Inches(2.10), "일정 조율 부담 감소"),
        (Inches(6.95), Inches(2.10), "변경 대응 속도 향상"),
        (LEFT, Inches(3.28), "실행력 향상"),
        (Inches(6.95), Inches(3.28), "기존 도구 대비 높은 실사용성 확보"),
    ]
    for x, y, text in items:
        add_text(slide, "•", x, y, Inches(0.20), Inches(0.18), size=BODY_SIZE, bold=True)
        add_text(slide, text, x + Inches(0.28), y - Inches(0.02), Inches(4.8), Inches(0.22), size=BODY_SIZE)
    add_rule(slide, LEFT, Inches(4.85), CONTENT_W)
    add_text(slide, "사람들은 기록 도구보다 조율 도구를 원한다.", LEFT, Inches(5.28), Inches(9.6), Inches(0.30), size=SUBTITLE_SIZE, bold=True)
    add_text(slide, "타임테이블은 그 공백을 메우는 AI 일정 비서 서비스다.", LEFT, Inches(5.82), Inches(10.5), Inches(0.30), size=SUBTITLE_SIZE, bold=True)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.title = "타임테이블 발표자료"
    prs.core_properties.subject = "AI 기반 일정 관리 비서 서비스"

    cover_slide(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)
    slide_8(prs)
    slide_9(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)
    slide_13(prs)
    slide_14(prs)
    slide_15(prs)
    slide_16(prs)
    slide_17(prs)
    slide_18(prs)
    slide_19(prs)
    slide_20(prs)
    slide_21(prs)
    slide_22(prs)
    slide_23(prs)
    slide_24(prs)
    slide_25(prs)
    slide_26(prs)
    slide_27(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


if __name__ == "__main__":
    print(build())
