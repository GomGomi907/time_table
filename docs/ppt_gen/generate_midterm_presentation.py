from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUTPUT_PATH = ROOT / "docs" / "Project_Midterm_Presentation_SpringBoot.pptx"
MOCKUP_IMAGE = ROOT / "docs" / "ppt_gen" / "assets" / "today_workspace_mockup.png"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(0xFF, 0xFF, 0xFF)
LAYER = RGBColor(0xF4, 0xF4, 0xF4)
NAVY = RGBColor(0x16, 0x16, 0x16)
GRAY = RGBColor(0x52, 0x52, 0x52)
DIVIDER = RGBColor(0xC6, 0xC6, 0xC6)
TEAL = RGBColor(0x0F, 0x62, 0xFE)
AMBER = RGBColor(0x00, 0x43, 0xCE)


def solid_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=16,
    bold=False,
    color=NAVY,
    align=PP_ALIGN.LEFT,
    font="Malgun Gothic",
    italic=False,
    margin=0.02,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_lines(
    slide,
    lines,
    x,
    y,
    w,
    h,
    *,
    size=14,
    color=NAVY,
    space_after=6,
    font="Malgun Gothic",
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        p.space_after = Pt(space_after)
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.color.rgb = color
    return box


def add_outline_box(slide, x, y, w, h, *, line_color=DIVIDER, line_width=1.0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LAYER
    shape.line.color.rgb = line_color
    shape.line.width = Pt(line_width)
    return shape


def add_accent_rule(slide, x, y, w, *, color=TEAL):
    rule = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        x,
        y,
        w,
        Inches(0.05),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = color
    rule.line.fill.background()
    return rule


def add_rule(slide, x, y, w, h, *, color=DIVIDER):
    rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    rule.fill.solid()
    rule.fill.fore_color.rgb = color
    rule.line.fill.background()
    return rule


def add_section_title(slide, title, subtitle):
    add_text(slide, "SPRING BOOT PROJECT", Inches(0.72), Inches(0.42), Inches(2.4), Inches(0.24), size=10, bold=True, color=TEAL)
    add_text(slide, title, Inches(0.72), Inches(0.82), Inches(8.0), Inches(0.58), size=28, bold=True)
    add_text(slide, subtitle, Inches(0.72), Inches(1.44), Inches(8.7), Inches(0.42), size=15, color=GRAY)


def add_page_number(slide, num):
    add_text(slide, f"{num:02d}", Inches(12.5), Inches(7.0), Inches(0.4), Inches(0.2), size=10, color=GRAY, align=PP_ALIGN.RIGHT)


def add_info_card(slide, x, y, w, h, title, body, *, accent=TEAL):
    add_outline_box(slide, x, y, w, h, line_color=NAVY, line_width=0.9)
    add_accent_rule(slide, x, y, w, color=accent)
    add_text(slide, title, x + Inches(0.18), y + Inches(0.16), w - Inches(0.36), Inches(0.3), size=16, bold=True)
    add_text(slide, body, x + Inches(0.18), y + Inches(0.52), w - Inches(0.36), h - Inches(0.7), size=13, color=GRAY)


def add_numbered_point(slide, num, title, body, x, y, w):
    add_text(slide, f"{num:02d}", x, y, Inches(0.4), Inches(0.2), size=12, bold=True, color=TEAL)
    add_rule(slide, x, y + Inches(0.28), Inches(0.28), Inches(0.03), color=TEAL)
    add_text(slide, title, x + Inches(0.52), y - Inches(0.04), w - Inches(0.56), Inches(0.24), size=17, bold=True)
    add_text(slide, body, x + Inches(0.52), y + Inches(0.20), w - Inches(0.56), Inches(0.40), size=12.5, color=GRAY)


def add_chip_row(slide, labels, x, y, w, chip_h=0.34, gap=0.12, color=NAVY):
    chip_w = (w - Inches(gap) * (len(labels) - 1)) / len(labels)
    for idx, label in enumerate(labels):
        cx = x + idx * (chip_w + Inches(gap))
        add_outline_box(slide, cx, y, chip_w, Inches(chip_h), line_color=color, line_width=0.8)
        add_text(slide, label, cx + Inches(0.06), y + Inches(0.06), chip_w - Inches(0.12), Inches(chip_h - 0.12), size=11, color=color, align=PP_ALIGN.CENTER)


def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide)
    add_text(slide, "AI Schedule Operator", Inches(0.82), Inches(1.05), Inches(5.2), Inches(0.6), size=30, bold=True)
    add_text(
        slide,
        "AI가 일정 변경 상황에서 다음 행동을 조율 및 제안하는 Spring Boot 기반 개인 스케줄 보조 시스템",
        Inches(0.82), Inches(1.78), Inches(5.9), Inches(0.7), size=16, color=GRAY
    )
    add_text(slide, "Project Midterm Presentation", Inches(0.82), Inches(2.75), Inches(2.8), Inches(0.28), size=11, bold=True, color=TEAL)
    add_text(slide, "소프트웨어공학 / Spring Boot 수업 프로젝트", Inches(0.82), Inches(6.15), Inches(3.8), Inches(0.25), size=12, color=GRAY)
    add_text(slide, "발표자: 이근산", Inches(0.82), Inches(6.45), Inches(2.0), Inches(0.25), size=12, color=GRAY)

    add_text(slide, "핵심 개념", Inches(7.45), Inches(1.35), Inches(1.3), Inches(0.22), size=14, bold=True, color=TEAL)
    add_accent_rule(slide, Inches(7.45), Inches(1.72), Inches(4.55), color=TEAL)
    add_text(slide, "입력", Inches(7.45), Inches(1.95), Inches(0.8), Inches(0.2), size=11, bold=True, color=TEAL)
    add_text(slide, "Goal  ·  Event  ·  Inbox", Inches(7.45), Inches(2.18), Inches(4.8), Inches(0.34), size=20, bold=True)
    add_text(slide, "목표, 외부 일정, 미배정 작업", Inches(7.45), Inches(2.56), Inches(4.6), Inches(0.22), size=12, color=GRAY)
    add_text(slide, "↓", Inches(9.4), Inches(3.0), Inches(0.3), Inches(0.2), size=22, color=GRAY, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(8.05), Inches(3.42), Inches(3.2), Inches(0.03), color=AMBER)
    add_text(slide, "AI Operator", Inches(8.05), Inches(3.62), Inches(3.4), Inches(0.36), size=24, bold=True)
    add_text(slide, "오늘 계획 생성 · 재조정안 제안", Inches(8.05), Inches(4.02), Inches(4.0), Inches(0.24), size=13, color=GRAY)
    add_text(slide, "↓", Inches(9.4), Inches(4.55), Inches(0.3), Inches(0.2), size=22, color=GRAY, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(7.45), Inches(4.98), Inches(4.55), Inches(0.03), color=TEAL)
    add_text(slide, "Today Workspace", Inches(7.45), Inches(5.2), Inches(4.2), Inches(0.3), size=22, bold=True)
    add_text(slide, "실행 화면에서 지금 해야 할 일과 다음 행동을 운영", Inches(7.45), Inches(5.58), Inches(4.7), Inches(0.26), size=13, color=GRAY)
    add_page_number(slide, 1)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide)
    add_section_title(slide, "왜 새로운 일정 앱이 필요한가", "일정 관리의 핵심 문제는 계획을 적는 것이 아니라, 흔들린 계획을 다시 운영하는 것이다.")
    add_numbered_point(slide, 1, "재조정 피로", "일정이 조금만 밀려도 사용자가 직접 다시 짜야 한다.", Inches(0.85), Inches(2.18), Inches(4.8))
    add_numbered_point(slide, 2, "도구 분리", "Goal, Todo, Calendar가 분리돼 오늘의 현실적인 계획이 잘 안 나온다.", Inches(0.85), Inches(3.25), Inches(4.8))
    add_numbered_point(slide, 3, "판단 지연", "실행 중 다음 행동을 빨리 결정하지 못해 집중이 더 깨진다.", Inches(0.85), Inches(4.32), Inches(4.8))

    add_text(slide, "기록 앱 vs 운영 앱", Inches(7.18), Inches(2.18), Inches(2.2), Inches(0.24), size=14, bold=True, color=TEAL)
    add_accent_rule(slide, Inches(7.18), Inches(2.56), Inches(5.05), color=TEAL)
    add_text(slide, "기존 앱", Inches(7.35), Inches(2.88), Inches(1.2), Inches(0.2), size=15, bold=True)
    add_text(slide, "제안 시스템", Inches(10.02), Inches(2.88), Inches(1.5), Inches(0.2), size=15, bold=True, color=TEAL)
    add_lines(slide, ["일정 기록 중심", "변경은 수동 처리", "오늘 실행 판단이 약함"], Inches(7.35), Inches(3.25), Inches(1.95), Inches(1.2), size=13, color=GRAY)
    add_lines(slide, ["일정 운영 중심", "변경안 자동 제안", "다음 행동 결정 지원"], Inches(10.0), Inches(3.25), Inches(2.2), Inches(1.2), size=13, color=GRAY)
    add_rule(slide, Inches(9.55), Inches(2.96), Inches(0.02), Inches(1.8), color=GRAY)
    add_rule(slide, Inches(7.18), Inches(5.08), Inches(5.05), Inches(0.02), color=GRAY)
    add_text(slide, "핵심 질문", Inches(7.18), Inches(5.25), Inches(1.0), Inches(0.2), size=12, bold=True, color=AMBER)
    add_text(slide, "\"기록 이후의 운영\"을 도와주는 도구가 필요한가?", Inches(8.15), Inches(5.2), Inches(4.15), Inches(0.4), size=13, color=NAVY)
    add_page_number(slide, 2)


def slide_persona(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide)
    add_section_title(slide, "핵심 사용자와 사용 맥락", "중간발표에서는 사용자가 누구인지보다, 어떤 운영 문제를 반복해서 겪는지 보여주는 것이 중요하다.")

    cols = [
        (Inches(0.82), Inches(3.15), "직장인 사이드빌더", "퇴근 후 공부와 프로젝트를 병행\n저녁 계획이 자주 무너짐", TEAL),
        (Inches(4.38), Inches(3.0), "계획형 학생", "강의·과제·시험 준비를 동시에 관리\n주간 목표와 오늘 실행 연결이 약함", AMBER),
        (Inches(7.82), Inches(4.2), "취준/포트폴리오형 사용자", "외부 일정과 자기계발이 섞여 있음\n시간 배분보다 재조정이 더 어려움", TEAL),
    ]
    for x, w, title, body, color in cols:
        add_accent_rule(slide, x, Inches(2.05), w, color=color)
        add_text(slide, title, x, Inches(2.24), w, Inches(0.28), size=18, bold=True)
        add_text(slide, body, x, Inches(2.62), w, Inches(0.5), size=12, color=GRAY)

    add_rule(slide, Inches(0.82), Inches(3.52), Inches(11.45), Inches(0.02), color=GRAY)
    add_rule(slide, Inches(4.16), Inches(3.78), Inches(0.02), Inches(2.0), color=GRAY)
    add_rule(slide, Inches(7.62), Inches(3.78), Inches(0.02), Inches(2.0), color=GRAY)

    add_text(slide, "하루 구조 예시", Inches(0.82), Inches(3.8), Inches(1.5), Inches(0.22), size=14, bold=True)
    add_lines(
        slide,
        ["1. 오전: 고정 일정", "2. 낮: 분절된 가용 시간", "3. 저녁: 핵심 블록 확보"],
        Inches(0.82),
        Inches(4.22),
        Inches(3.05),
        Inches(0.9),
        size=12,
        color=NAVY,
        space_after=2,
    )
    add_text(slide, "통근과 수업/업무로 가용 시간이 파편화된다.", Inches(0.82), Inches(5.28), Inches(3.1), Inches(0.28), size=11.5, color=GRAY)

    add_text(slide, "주간 구조 예시", Inches(4.4), Inches(3.8), Inches(1.5), Inches(0.22), size=14, bold=True)
    add_lines(
        slide,
        ["1. 평일: 밀도 높음", "2. 주말: 회복과 보충", "3. 일요일: 재정비"],
        Inches(4.4),
        Inches(4.22),
        Inches(2.8),
        Inches(0.9),
        size=12,
        color=NAVY,
        space_after=2,
    )
    add_text(slide, "주간 리듬이 있어 루틴과 보호 시간이 중요하다.", Inches(4.4), Inches(5.28), Inches(2.9), Inches(0.28), size=11.5, color=GRAY)

    add_text(slide, "대표 실패 시나리오", Inches(7.86), Inches(3.8), Inches(1.8), Inches(0.22), size=14, bold=True, color=AMBER)
    add_lines(
        slide,
        [
            "1. 회의 연장으로 저녁 계획 붕괴",
            "2. 현재 작업 초과로 뒤 블록 취소",
            "3. 급한 일 유입으로 원래 목표 지연",
            "4. 캘린더 누락으로 현실과 계획 불일치",
            "5. 후보는 많지만 오늘 할 일 미결정",
        ],
        Inches(7.86),
        Inches(4.18),
        Inches(4.0),
        Inches(1.2),
        size=11,
        color=GRAY,
        space_after=2,
    )
    add_page_number(slide, 3)


def slide_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide)
    add_section_title(slide, "제안 솔루션", "목표, 루틴, 외부 일정, 인박스를 한데 모아 AI가 오늘 계획과 재조정안을 제안한다.")
    steps = [
        ("1", "입력", "Goal / Inbox"),
        ("2", "제약", "Routine / Calendar"),
        ("3", "생성", "오늘 계획 초안"),
        ("4", "실행", "Today Workspace"),
        ("5", "개입", "재조정 제안"),
    ]
    for idx, (num, title, desc) in enumerate(steps):
        x = Inches(0.95) + Inches(2.45) * idx
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.6), Inches(2.25), Inches(0.48), Inches(0.48))
        circ.fill.solid()
        circ.fill.fore_color.rgb = TEAL
        circ.line.fill.background()
        add_text(slide, num, x + Inches(0.6), Inches(2.30), Inches(0.48), Inches(0.16), size=12, bold=True, color=BG, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(0.2), Inches(2.95), Inches(1.3), Inches(0.22), size=16, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, desc, x + Inches(0.05), Inches(3.32), Inches(1.6), Inches(0.34), size=12, color=GRAY, align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_text(slide, "→", x + Inches(1.78), Inches(2.78), Inches(0.3), Inches(0.2), size=20, color=GRAY, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(1.2), Inches(4.9), Inches(10.9), Inches(0.02), color=GRAY)
    add_text(slide, "핵심 원칙", Inches(1.2), Inches(5.1), Inches(1.1), Inches(0.22), size=14, bold=True, color=TEAL)
    add_text(slide, "AI는 승인 없는 자동 재조정을 하지 않고, 항상 이유와 영향 범위를 함께 설명한다.", Inches(2.55), Inches(5.06), Inches(8.9), Inches(0.28), size=14, color=NAVY)
    add_page_number(slide, 4)


def slide_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide)
    add_section_title(slide, "핵심 기능", "중간발표 범위에서는 복잡한 플랫폼보다, 일정 운영에 필요한 최소 기능에 집중한다.")
    add_rule(slide, Inches(6.34), Inches(2.08), Inches(0.02), Inches(3.7), color=GRAY)
    add_numbered_point(slide, 1, "목표 계획", "장기·주간 Goal을 만들고 오늘 후보 작업으로 연결", Inches(0.95), Inches(2.18), Inches(4.95))
    add_numbered_point(slide, 2, "루틴/외부 일정 수집", "RoutineTemplate와 Google Calendar로 보호 시간을 계산", Inches(6.65), Inches(2.18), Inches(5.0))
    add_rule(slide, Inches(0.95), Inches(3.86), Inches(5.0), Inches(0.02), color=GRAY)
    add_rule(slide, Inches(6.65), Inches(3.86), Inches(5.0), Inches(0.02), color=GRAY)
    add_numbered_point(slide, 3, "오늘 워크스페이스", "FocusSession, 타임라인, 인박스, 목표를 한 화면에 통합", Inches(0.95), Inches(4.18), Inches(4.95))
    add_numbered_point(slide, 4, "AI 개입", "지연·충돌·빈 시간 발생 시 재조정안을 제안 후 승인", Inches(6.65), Inches(4.18), Inches(5.0))
    add_text(slide, "범위 통제", Inches(0.92), Inches(6.28), Inches(0.9), Inches(0.2), size=12, bold=True, color=TEAL)
    add_text(slide, "팀 협업, 마켓플레이스, 통계 리포트는 제외하고 Spring Boot 구현 가능 범위를 우선한다.", Inches(1.95), Inches(6.24), Inches(9.9), Inches(0.24), size=12, color=GRAY)
    add_page_number(slide, 5)


def slide_uiux(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide)
    add_section_title(slide, "UI/UX 방향과 메인 화면 목업", "Simple is Best, 인지 부하 최소화, F-패턴 배치를 기준으로 메인 화면을 설계했다.")
    add_text(slide, "핵심 UX 원칙", Inches(0.9), Inches(2.2), Inches(1.5), Inches(0.22), size=14, bold=True, color=TEAL)
    add_numbered_point(slide, 1, "Execution First", "현재 해야 할 일을 가장 먼저 읽히게 배치", Inches(0.92), Inches(2.72), Inches(2.6))
    add_numbered_point(slide, 2, "Explainable AI", "개입 이유와 영향 범위를 짧게 설명", Inches(0.92), Inches(3.92), Inches(2.6))
    add_numbered_point(slide, 3, "Recovery UI", "문제보다 다음 행동을 먼저 보여줌", Inches(0.92), Inches(5.12), Inches(2.6))
    add_rule(slide, Inches(4.08), Inches(2.1), Inches(0.02), Inches(4.55), color=GRAY)
    if MOCKUP_IMAGE.exists():
        slide.shapes.add_picture(str(MOCKUP_IMAGE), Inches(4.45), Inches(2.2), width=Inches(7.45), height=Inches(4.18))
    add_rule(slide, Inches(4.45), Inches(6.48), Inches(7.45), Inches(0.02), color=GRAY)
    add_text(slide, "① 포커스 보드   ② AI 개입 카드   ③ 오늘 타임라인", Inches(4.55), Inches(6.54), Inches(7.2), Inches(0.2), size=11, color=GRAY, align=PP_ALIGN.CENTER)
    add_page_number(slide, 6)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide)
    add_section_title(slide, "Spring Boot 구현 구조", "단일 Spring Boot 앱 안에서 인증, 일정 엔진, 외부 연동, AI adapter를 분리해 구현한다.")
    blocks = [
        (Inches(0.92), Inches(2.35), "Web", ["Thymeleaf", "Today Workspace", "설정/온보딩"], TEAL),
        (Inches(3.95), Inches(2.65), "Spring Boot Core", ["Controller", "Service", "Schedule Engine", "AI Adapter"], AMBER),
        (Inches(7.15), Inches(2.4), "DB", ["Spring Data JPA", "MySQL", "Goal / Event / Block"], AMBER),
        (Inches(10.15), Inches(2.15), "External", ["Google Calendar API", "OAuth2 Login", "LLM API"], TEAL),
    ]
    for idx, (x, w, title, body_lines, color) in enumerate(blocks):
        add_accent_rule(slide, x, Inches(2.24), w, color=color)
        add_text(slide, title, x, Inches(2.48), w, Inches(0.28), size=18, bold=True)
        add_lines(slide, body_lines, x, Inches(2.92), w, Inches(1.45), size=13, color=GRAY, space_after=2)
        if idx < len(blocks) - 1:
            add_text(slide, "→", x + w + Inches(0.08), Inches(3.28), Inches(0.3), Inches(0.2), size=20, color=GRAY, align=PP_ALIGN.CENTER)
    add_rule(slide, Inches(0.92), Inches(5.22), Inches(11.3), Inches(0.02), color=GRAY)
    add_text(slide, "구현 포인트", Inches(0.92), Inches(5.42), Inches(1.0), Inches(0.2), size=12, bold=True, color=TEAL)
    add_text(slide, "AI가 없어도 기본 일정 생성은 작동하고, AI가 붙으면 재조정안 품질이 좋아지는 구조로 설계한다.", Inches(2.0), Inches(5.38), Inches(9.8), Inches(0.24), size=12, color=GRAY)
    add_page_number(slide, 7)


def slide_progress(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide)
    add_section_title(slide, "중간 진행 상황과 남은 일정", "기획과 UX 기준은 정리되었고, 다음 단계는 Spring Boot 구현과 데모 흐름 연결이다.")
    add_text(slide, "완료된 것", Inches(0.95), Inches(2.18), Inches(1.4), Inches(0.24), size=18, bold=True)
    add_accent_rule(slide, Inches(0.95), Inches(2.56), Inches(1.15), color=TEAL)
    add_lines(
        slide,
        ["문제 정의", "페르소나 정리", "기능/페이지 명세", "UI/UX 설계", "메인 화면 목업"],
        Inches(0.95),
        Inches(2.88),
        Inches(4.9),
        Inches(1.5),
        size=13,
        color=GRAY,
        space_after=2,
    )

    add_rule(slide, Inches(6.3), Inches(2.15), Inches(0.02), Inches(2.1), color=GRAY)
    add_text(slide, "다음 구현 항목", Inches(6.62), Inches(2.18), Inches(2.45), Inches(0.42), size=17, bold=True)
    add_accent_rule(slide, Inches(6.62), Inches(2.56), Inches(1.35), color=AMBER)
    add_lines(
        slide,
        ["OAuth2 로그인", "JPA 엔티티 설계", "Google Calendar 연동", "Today Workspace CRUD", "간단한 AI 개입 API"],
        Inches(6.62),
        Inches(2.88),
        Inches(5.0),
        Inches(1.5),
        size=13,
        color=GRAY,
        space_after=2,
    )

    add_text(slide, "데모 목표", Inches(0.95), Inches(4.72), Inches(0.9), Inches(0.2), size=12, bold=True, color=TEAL)
    add_text(slide, "로그인 → 연동 → 오늘 계획 생성 → 일정 지연 → AI 개입 승인", Inches(1.95), Inches(4.69), Inches(9.8), Inches(0.22), size=12.5, color=NAVY)
    add_rule(slide, Inches(1.15), Inches(5.5), Inches(10.85), Inches(0.02), color=GRAY)

    weeks = [
        (Inches(1.35), "1주차", "프로젝트 세팅\n로그인 / 엔티티"),
        (Inches(4.0), "2주차", "Calendar 연동\nRoutine / Inbox CRUD"),
        (Inches(6.75), "3주차", "오늘 계획 생성\nWorkspace 조회"),
        (Inches(9.45), "4주차", "AI 개입 API\n최종 데모 정리"),
    ]
    for idx, (x, week, body) in enumerate(weeks):
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, Inches(5.39), Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL if idx < 3 else AMBER
        dot.line.fill.background()
        add_text(slide, week, x - Inches(0.02), Inches(5.72), Inches(1.2), Inches(0.22), size=14, bold=True, color=TEAL if idx == 3 else NAVY)
        add_text(slide, body, x - Inches(0.02), Inches(6.02), Inches(2.1), Inches(0.46), size=11.5, color=GRAY)
    add_page_number(slide, 8)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.title = "Project Midterm Presentation - AI Schedule Operator"
    prs.core_properties.subject = "Spring Boot class project"

    slide_title(prs)
    slide_problem(prs)
    slide_persona(prs)
    slide_solution(prs)
    slide_features(prs)
    slide_uiux(prs)
    slide_architecture(prs)
    slide_progress(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


if __name__ == "__main__":
    print(build())
