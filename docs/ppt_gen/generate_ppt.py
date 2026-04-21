import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def set_font(run, font_name="맑은 고딕", size=Pt(18), bold=False, color=RGBColor(0x33, 0x33, 0x33)):
    font = run.font
    font.name = font_name
    font.size = size
    font.bold = bold
    font.color.rgb = color

def add_slide_with_title(prs, layout_idx, title_text):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    title = slide.shapes.title
    title.text = title_text
    # 제목 폰트 설정
    for paragraph in title.text_frame.paragraphs:
        for run in paragraph.runs:
            set_font(run, size=Pt(32), bold=True, color=RGBColor(0x2B, 0x6C, 0xB0)) # 포인트 컬러 사용
    return slide

def add_bullet_points(tf, points, level=0, size=Pt(18)):
    for p_text in points:
        p = tf.add_paragraph()
        p.text = p_text
        p.level = level
        for run in p.runs:
            set_font(run, size=size)

def main():
    prs = Presentation()
    
    # 기본 슬라이드 레이아웃
    TITLE_SLIDE_LAYOUT = 0
    BULLET_SLIDE_LAYOUT = 1
    BLANK_SLIDE_LAYOUT = 6

    # ------------------
    # Slide 1: 타이틀
    # ------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE_LAYOUT])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]

    title.text = "타임테이블 프로젝트 기획"
    for run in title.text_frame.paragraphs[0].runs:
        set_font(run, size=Pt(44), bold=True, color=RGBColor(0x2B, 0x6C, 0xB0))
        
    subtitle.text = "LLM 에이전트 기반 지능형 일정 관리\n\n단국대학교 컴퓨터공학과\n이근산"
    for p in subtitle.text_frame.paragraphs:
        for run in p.runs:
            set_font(run, size=Pt(20), color=RGBColor(0x55, 0x55, 0x55))
            
    # [아이콘 배치 자리 안내 추가]
    txBox = slide1.shapes.add_textbox(Inches(3), Inches(1), Inches(4), Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "[아이콘: 캘린더와 톱니바퀴 결합]"
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs: set_font(run, size=Pt(16), color=RGBColor(0x77, 0x77, 0x77))

    # ------------------
    # Slide 2: 문제 제기 (Problem)
    # ------------------
    slide2 = add_slide_with_title(prs, BULLET_SLIDE_LAYOUT, "문제 제기 (Problem Definition)")
    body2 = slide2.shapes.placeholders[1]
    tf2 = body2.text_frame
    tf2.text = ""
    add_bullet_points(tf2, [
        "일정 관리의 본질적 어려움",
        " - 계획을 '세우는 것'보다 '수정하고 유지하는 것'에 큰 에너지 소모",
        " - 실행 기능(Executive Function) 저하 시 일정 관리가 가장 먼저 무너짐",
        "기존 일정 앱의 한계",
        " - 단순히 예쁘게 '기록'하고 '조회'하는 수동적 도구",
        " - 일정 지연/취소 시 사용자가 직접 재배치해야 하는 '관리 피로' 발생"
    ])
    txBox = slide2.shapes.add_textbox(Inches(7), Inches(2), Inches(2.5), Inches(3))
    txBox.text_frame.text = "[아이콘]\n마음대로 안되는 일정\n\n'계획 수립 << 계획 유지'"

    # ------------------
    # Slide 3: 솔루션 (Solution)
    # ------------------
    slide3 = add_slide_with_title(prs, BULLET_SLIDE_LAYOUT, "프로젝트 목적 및 솔루션")
    tf3 = slide3.shapes.placeholders[1].text_frame
    tf3.text = ""
    add_bullet_points(tf3, [
        "목적: 부족한 인지적, 실행적 에너지를 보조하는 외부 시스템",
        "솔루션: \"일정을 함께 운영하는 비서형 앱\"",
        " 1. 패턴 보존 [방패 아이콘] : 기본 주간 패턴(수면, 식사 등) 최우선 보호",
        " 2. 부하 감소 [번개 아이콘] : 자연어 요청으로 5분 단위 조율안 제시",
        " 3. 자동 보정 [새로고침 아이콘] : 일정 종료 후 상태 입력에 따른 자동 재배치"
    ])

    # ------------------
    # Slide 4: 핵심 시나리오
    # ------------------
    slide4 = add_slide_with_title(prs, BULLET_SLIDE_LAYOUT, "핵심 타겟 및 사용자 시나리오")
    tf4 = slide4.shapes.placeholders[1].text_frame
    tf4.text = ""
    add_bullet_points(tf4, [
        "메인 타겟: 잦은 변동성에 취약하거나 일정 관리에 피로를 느끼는 학생",
        "핵심 시나리오:",
        " [퍼즐 아이콘] 1. 온보딩 : 구글 로그인 후 요일별 고정 루틴 입력 (기준점 형성)",
        " [말풍선 아이콘] 2. 조율 요청 : \"저녁에 피곤한데 헬스는 내일로 미루고 과제 배치해줘\"",
        " [톱니바퀴 아이콘] 3. 자동 재배치 : 지연 예상 시 다음 스케줄을 밀어내고 승인 요청"
    ])

    # ------------------
    # Slide 5: 주요 기능 1
    # ------------------
    slide5 = add_slide_with_title(prs, BULLET_SLIDE_LAYOUT, "기능 1: 데이터 통합 및 직관적 대시보드")
    tf5 = slide5.shapes.placeholders[1].text_frame
    tf5.text = ""
    add_bullet_points(tf5, [
        "인지 부하를 최소화하는 대시보드",
        " - 복잡한 뷰 대신 \"지금 당장 해야 하는 일\" 1개만 최상단 강조",
        " - 오늘 일정 타임라인을 시간순으로 직관적 나열",
        "통합 연동 [구글 로고 아이콘]",
        " - Google Calendar (외부일정) + Google Tasks (할일) 완벽 통합",
        " - 복잡성 제거 -> 직관적 행동 유도"
    ])

    # ------------------
    # Slide 6: 주요 기능 2
    # ------------------
    slide6 = add_slide_with_title(prs, BULLET_SLIDE_LAYOUT, "기능 2: LLM 에이전트와 규칙 검증기")
    tf6 = slide6.shapes.placeholders[1].text_frame
    tf6.text = ""
    add_bullet_points(tf6, [
        "자연어 기반 조율",
        " - 사용자는 상황만 설명, 복잡한 테트리스는 시스템이 수행",
        "안전한 조율 구조 파이프라인",
        " 1. [사용자] 자연어 명령",
        " 2. [로봇 아이콘] LLM의 제안안 (JSON)",
        " 3. [거름망 아이콘] 시스템 규칙 검증기 (고정일정 보호, 중첩 불가 확인)",
        " 4. [캘린더 아이콘] 사용자 미리보기 및 최종 승인"
    ])

    # ------------------
    # Slide 7: 주요 기능 3
    # ------------------
    slide7 = add_slide_with_title(prs, BULLET_SLIDE_LAYOUT, "기능 3: 현실 기반 상태 추적과 동적 재조정")
    tf7 = slide7.shapes.placeholders[1].text_frame
    tf7.text = ""
    add_bullet_points(tf7, [
        "현실과의 동기화: 일정 종료 시점에서 팝업/알림으로 실제 상태 확인",
        "5가지 상태 전이 모델:",
        " [✅ 완료] 그대로 유지",
        " [⏩ 조기 완료] 빈 슬롯 발생 -> 다음 작업 당기기",
        " [⏳ 연장] 뒤 일정을 순차적으로 밀기",
        " [↪️ 미루기] 다음 가능 슬롯 재탐색",
        " [❌ 취소] 삭제"
    ])

    # ------------------
    # Slide 8: 
    # ------------------
    slide8 = add_slide_with_title(prs, BULLET_SLIDE_LAYOUT, "시스템 아키텍처 및 기술 스택")
    tf8 = slide8.shapes.placeholders[1].text_frame
    tf8.text = ""
    add_bullet_points(tf8, [
        "[Client]",
        " - Web UI (Thymeleaf / JS 기반 렌더링)",
        "[Server]",
        " - Spring Boot 3.x",
        " - Spring Security (OAuth2 로그인)",
        " - LLM Adapter 계층",
        "[Data]",
        " - PostgreSQL (시간/상태 무결성)",
        " - Google Calendar / Tasks API 연동"
    ])

    # ------------------
    # Slide 9: 
    # ------------------
    slide9 = add_slide_with_title(prs, BULLET_SLIDE_LAYOUT, "8주 개발 로드맵 (1인 개발 기준)")
    tf9 = slide9.shapes.placeholders[1].text_frame
    tf9.text = ""
    add_bullet_points(tf9, [
        "Week 1~3: [기반/건물 아이콘] 기반 구축 및 온보딩",
        " - Spring Boot 세팅, Google OAuth, 루틴 등록 API",
        "Week 4~5: [링크 아이콘] 데이터 연동",
        " - 일정/태스크 CRUD, 대시보드, Google Sync",
        "Week 6~7: [뇌 아이콘] 지능형 조율",
        " - LLM 프롬프트, 제안안 처리, 검증기",
        "Week 8: [깃발 아이콘] 상태 반영 및 마무리",
        " - 동적 재조정 로직 시연 리허설"
    ])

    # ------------------
    # Slide 10: 
    # ------------------
    slide10 = add_slide_with_title(prs, TITLE_SLIDE_LAYOUT, "기대 효과 및 마무리")
    title10 = slide10.shapes.title
    title10.text = "기대 효과"
    for run in title10.text_frame.paragraphs[0].runs:
         set_font(run, size=Pt(40), bold=True, color=RGBColor(0x2B, 0x6C, 0xB0))
         
    subtitle10 = slide10.placeholders[1]
    subtitle10.text = "\n[차 마시는 사람 아이콘]\n\n일정 유지보수에 낭비되는 인지적 에너지를 획기적으로 절약\n\n\"계획이 틀어져도 포기하지 않는 지속 가능한 시간 관리\""
    for p in subtitle10.text_frame.paragraphs:
        for run in p.runs:
            set_font(run, size=Pt(24), bold=True, color=RGBColor(0x33, 0x33, 0x33))
            
    # 파일 저장
    output_path = "TimeTable_Project_Proposal.pptx"
    prs.save(output_path)
    print(f"PPT 생성 완료: {output_path}")

if __name__ == "__main__":
    main()
